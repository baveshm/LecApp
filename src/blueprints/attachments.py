"""Attachments blueprint (Task 10)

Extracted from monolithic routes in `app.py`:
 - GET /api/recordings/<recording_id>/files
 - POST /api/recordings/<recording_id>/files
 - GET /api/recordings/<recording_id>/files/<file_id>/preview
 - DELETE /api/recordings/<recording_id>/files/<file_id>

Routes maintain identical behavior & response payloads for test stability.
"""
from __future__ import annotations

import os
import mimetypes
from datetime import datetime

from flask import Blueprint, jsonify, request, send_file, make_response, current_app
from flask_login import login_required, current_user
from werkzeug.utils import secure_filename

from src.extensions import db
from src.models import Recording, Attachment, SystemSetting

attachments_bp = Blueprint('attachments', __name__)


def _own_recording_or_404(recording_id: int):
    recording = db.session.get(Recording, recording_id)
    if not recording:
        return None, (jsonify({'error': 'Recording not found'}), 404)
    if recording.user_id != current_user.id:
        return None, (jsonify({'error': 'Unauthorized'}), 403)
    return recording, None


@attachments_bp.route('/api/recordings/<int:recording_id>/files', methods=['GET'])
@login_required
def get_recording_files(recording_id: int):
    try:
        recording, error = _own_recording_or_404(recording_id)
        if error:
            return error
        attachments = (Attachment.query
                        .filter_by(recording_id=recording_id)
                        .order_by(Attachment.uploaded_at.desc())
                        .all())
        return jsonify([a.to_dict() for a in attachments])
    except Exception as e:
        current_app.logger.error(f"Error getting files for recording {recording_id}: {e}")
        return jsonify({'error': str(e)}), 500


@attachments_bp.route('/api/recordings/<int:recording_id>/files', methods=['POST'])
@login_required
def upload_recording_file(recording_id: int):
    try:
        recording, error = _own_recording_or_404(recording_id)
        if error:
            return error
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        file = request.files['file']
        if not file.filename:
            return jsonify({'error': 'No file selected'}), 400
        allowed_extensions = {
            '.pdf', '.doc', '.docx', '.txt', '.md', '.rtf',
            '.ppt', '.pptx', '.xls', '.xlsx', '.csv', '.json',
            '.png', '.jpg', '.jpeg', '.heic', '.webp'
        }
        filename = file.filename
        file_ext = os.path.splitext(filename)[1].lower()
        if file_ext not in allowed_extensions:
            return jsonify({'error': f'File type {file_ext} not allowed. Allowed types: {", ".join(allowed_extensions)}'}), 400
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        max_file_size_mb = SystemSetting.get_setting('max_file_size_mb', 250) or 250
        max_size_bytes = max_file_size_mb * 1024 * 1024
        if file_size > max_size_bytes:
            return jsonify({'error': f'File too large. Maximum size is {max_file_size_mb}MB'}), 413
        original_filename = filename
        safe_name = secure_filename(original_filename)
        from flask import current_app as app
        attachments_dir = os.path.join(app.config['UPLOAD_FOLDER'], 'attachments', str(recording_id))
        os.makedirs(attachments_dir, exist_ok=True)
        timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
        file_path = os.path.join(attachments_dir, f"{timestamp}_{safe_name}")
        file.save(file_path)
        current_app.logger.info(f"Saved attachment to {file_path}")
        mime_type, _ = mimetypes.guess_type(original_filename)
        file_type = file_ext.lstrip('.')
        attachment = Attachment(
            recording_id=recording_id,
            user_id=current_user.id,
            file_path=file_path,
            original_filename=original_filename,
            file_size=file_size,
            file_type=file_type,
            mime_type=mime_type
        )
        db.session.add(attachment)
        db.session.commit()
        current_app.logger.info(f"Created attachment record {attachment.id} for recording {recording_id}")
        return jsonify(attachment.to_dict()), 201
    except Exception as e:
        db.session.rollback()
        current_app.logger.error(f"Error uploading file to recording {recording_id}: {e}")
        return jsonify({'error': str(e)}), 500


@attachments_bp.route('/api/recordings/<int:recording_id>/files/<int:file_id>/preview', methods=['GET'])
@login_required
def preview_recording_file(recording_id: int, file_id: int):
    try:
        recording, error = _own_recording_or_404(recording_id)
        if error:
            return error
        attachment = Attachment.query.filter_by(id=file_id, recording_id=recording_id).first()
        if not attachment:
            return jsonify({'error': 'File not found'}), 404
        if not os.path.exists(attachment.file_path):
            current_app.logger.error(f"Attachment file missing: {attachment.file_path}")
            return jsonify({'error': 'File not found on server'}), 404
        file_type = (attachment.file_type or '').lower()
        # Images
        if file_type in ['png', 'jpg', 'jpeg', 'heic', 'webp']:
            return send_file(attachment.file_path, mimetype=attachment.mime_type or 'image/jpeg')
        # PDF (first 2 pages)
        if file_type == 'pdf':
            try:
                from PyPDF2 import PdfReader, PdfWriter
                from io import BytesIO
                reader = PdfReader(attachment.file_path)
                writer = PdfWriter()
                for i in range(min(2, len(reader.pages))):
                    writer.add_page(reader.pages[i])
                preview = BytesIO()
                writer.write(preview)
                preview.seek(0)
                return send_file(preview, mimetype='application/pdf', as_attachment=False, download_name=f"preview_{attachment.original_filename}")
            except ImportError:
                current_app.logger.warning("PyPDF2 not available, returning full PDF")
                return send_file(attachment.file_path, mimetype='application/pdf')
            except Exception as e:
                current_app.logger.error(f"Error generating PDF preview: {e}")
                return jsonify({'error': 'Failed to generate PDF preview'}), 500
        # Text-like
        if file_type in ['txt', 'md', 'rtf', 'json', 'csv']:
            try:
                with open(attachment.file_path, 'r', encoding='utf-8') as f:
                    content = f.read(2000)
                if len(content) == 2000:
                    content += "\n\n[Preview truncated - showing first 2000 characters]"
                resp = make_response(content)
                resp.headers['Content-Type'] = 'text/plain; charset=utf-8'
                return resp
            except Exception as e:
                current_app.logger.error(f"Error reading text file: {e}")
                return jsonify({'error': 'Failed to read file'}), 500
        # Word
        if file_type in ['doc', 'docx']:
            try:
                from docx import Document
                from io import BytesIO
                doc = Document(attachment.file_path)
                paragraphs = []
                count = 0
                max_chars = 2000
                for para in doc.paragraphs:
                    if count >= max_chars:
                        break
                    paragraphs.append(para.text)
                    count += len(para.text)
                preview_text = '\n\n'.join(paragraphs)
                if count >= max_chars:
                    preview_text += "\n\n[Preview truncated - showing approximately first 2 pages]"
                resp = make_response(preview_text)
                resp.headers['Content-Type'] = 'text/plain; charset=utf-8'
                return resp
            except ImportError:
                current_app.logger.warning("python-docx not available")
                return jsonify({'error': 'Document preview not available - library not installed'}), 501
            except Exception as e:
                current_app.logger.error(f"Error reading Word document: {e}")
                return jsonify({'error': 'Failed to read document'}), 500
        # Excel
        if file_type in ['xls', 'xlsx']:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(attachment.file_path, read_only=True)
                if wb.worksheets:
                    sheet = wb.active
                    rows = []
                    for i, row in enumerate(sheet.iter_rows(values_only=True)):
                        if i >= 2:
                            break
                        rows.append('\t'.join(str(c) if c is not None else '' for c in row))
                    preview_text = '\n'.join(rows) + "\n\n[Preview showing first 2 rows]"
                    resp = make_response(preview_text)
                    resp.headers['Content-Type'] = 'text/plain; charset=utf-8'
                    return resp
                else:
                    return jsonify({'error': 'No worksheets found in workbook'}), 500
            except ImportError:
                current_app.logger.warning("openpyxl not available")
                return jsonify({'error': 'Spreadsheet preview not available - library not installed'}), 501
            except Exception as e:
                current_app.logger.error(f"Error reading spreadsheet: {e}")
                return jsonify({'error': 'Failed to read spreadsheet'}), 500
        # PowerPoint
        if file_type in ['ppt', 'pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(attachment.file_path)
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    if i >= 2:
                        break
                    slide_text = f"=== Slide {i+1} ===\n"
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            slide_text += shape.text + '\n'
                    slides_text.append(slide_text)
                preview_text = '\n\n'.join(slides_text) + "\n\n[Preview showing first 2 slides]"
                resp = make_response(preview_text)
                resp.headers['Content-Type'] = 'text/plain; charset=utf-8'
                return resp
            except ImportError:
                current_app.logger.warning("python-pptx not available")
                return jsonify({'error': 'Presentation preview not available - library not installed'}), 501
            except Exception as e:
                current_app.logger.error(f"Error reading presentation: {e}")
                return jsonify({'error': 'Failed to read presentation'}), 500
        return jsonify({'error': f'Preview not available for {file_type} files'}), 501
    except Exception as e:
        current_app.logger.error(f"Error previewing file {file_id} for recording {recording_id}: {e}")
        return jsonify({'error': str(e)}), 500


@attachments_bp.route('/api/recordings/<int:recording_id>/files/<int:file_id>/download', methods=['GET'])
@login_required
def download_recording_file(recording_id: int, file_id: int):
    try:
        recording, error = _own_recording_or_404(recording_id)
        if error:
            return error
        attachment = Attachment.query.filter_by(id=file_id, recording_id=recording_id).first()
        if not attachment:
            return jsonify({'error': 'File not found'}), 404
        if not os.path.exists(attachment.file_path):
            current_app.logger.error(f"Attachment file missing: {attachment.file_path}")
            return jsonify({'error': 'File not found on server'}), 404
        return send_file(
            attachment.file_path,
            mimetype=attachment.mime_type or 'application/octet-stream',
            as_attachment=True,
            download_name=attachment.original_filename
        )
    except Exception as e:
        current_app.logger.error(f"Error downloading file {file_id} for recording {recording_id}: {e}")
        return jsonify({'error': str(e)}), 500


@attachments_bp.route('/api/recordings/<int:recording_id>/files/<int:file_id>/unlink', methods=['POST'])
@login_required
def unlink_recording_file(recording_id: int, file_id: int):
    try:
        recording, error = _own_recording_or_404(recording_id)
        if error:
            return error
        attachment = Attachment.query.filter_by(id=file_id, recording_id=recording_id).first()
        if not attachment:
            return jsonify({'error': 'File not found'}), 404
        # Remove the recording association but keep the file and attachment record
        attachment.recording_id = None
        db.session.commit()
        current_app.logger.info(f"Unlinked attachment {file_id} from recording {recording_id}")
        return jsonify({'success': True, 'message': 'File unlinked successfully'})
    except Exception as e:
        db.session.rollback()
        current_app.logger.error(f"Error unlinking file {file_id} from recording {recording_id}: {e}")
        return jsonify({'error': str(e)}), 500


@attachments_bp.route('/api/recordings/<int:recording_id>/files/<int:file_id>', methods=['DELETE'])
@login_required
def delete_recording_file(recording_id: int, file_id: int):
    try:
        recording, error = _own_recording_or_404(recording_id)
        if error:
            return error
        attachment = Attachment.query.filter_by(id=file_id, recording_id=recording_id).first()
        if not attachment:
            return jsonify({'error': 'File not found'}), 404
        try:
            if os.path.exists(attachment.file_path):
                os.remove(attachment.file_path)
                current_app.logger.info(f"Deleted attachment file: {attachment.file_path}")
        except Exception as e:
            current_app.logger.error(f"Error deleting attachment file {attachment.file_path}: {e}")
        db.session.delete(attachment)
        db.session.commit()
        current_app.logger.info(f"Deleted attachment {file_id} from recording {recording_id}")
        return jsonify({'success': True, 'message': 'File deleted successfully'})
    except Exception as e:
        db.session.rollback()
        current_app.logger.error(f"Error deleting file {file_id} from recording {recording_id}: {e}")
        return jsonify({'error': str(e)}), 500
